# from langchain.document_loaders import PyPDFium2Loader
import pandas as pd
import fitz
from docx import Document
from pptx import Presentation
import os
import json
# import  PyMuPDF



def read_docx(file_path):
    try:

        doc = Document(file_path)
        text = ""

     
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"

        return text

    except Exception as e:
        print(f"Error: {e}")
        return None



def read_pdf(pdf_path):
    try:
       
        pdf_document = fitz.open(pdf_path)
        text = ""

        # Iterate through each page in the PDF
        for page_number in range(pdf_document.page_count):
            # Get the page
            page = pdf_document[page_number]
            text += page.get_text()

        # Close the PDF document
        pdf_document.close()

        return text

    except Exception as e:
        print(f"Error: {e}")
        return None


def read_pptx(pptx_path):
    try:
       
        presentation = Presentation(pptx_path)
        text = ""

        # Iterate through each slide in the presentation
        for slide_number, slide in enumerate(presentation.slides, start=1):
            text += f"\nSlide {slide_number}:\n"

            # Iterate through each shape in the slide
            for shape_number, shape in enumerate(slide.shapes, start=1):
                if hasattr(shape, "text"):
                    text += f"  Shape {shape_number}: {shape.text}\n"

        return text

    except Exception as e:
        print(f"Error: {e}")
        return None




def read_excel(file_path):
    try:
        # Read the Excel file into a DataFrame
        df = pd.read_excel(file_path)
        print(df)

        # Optionally, return the DataFrame if you want to use it in your script
        return df

    except Exception as e:
        print(f"Error: {e}")
        return None
    
    
    
def read_csv(file_path):
    
    data = pd.read_csv(file_path)
    return data







def read_json(file_path):
    with open(file_path, 'r') as json_file:
            data = json.load(json_file)
    return data




    


def read_file(file_path):
    _, extension = os.path.splitext(file_path.lower())
    if extension == '.docx':
        return read_docx(file_path)
    elif extension == '.pdf':
        return read_pdf(file_path)
    elif extension == '.pptx':
        return read_pptx(file_path)
    elif extension == '.xlsx':
        return read_excel(file_path)
    elif extension == '.csv':
        return read_csv(file_path)
    elif extension == '.json':
        return read_json(file_path)
    else:
        print(f"Unsupported file extension: {extension}")
        return None
    
   
   
   
def read_files_in_folder(folder_path):
    
    contents = []
    for root, dirs, files in os.walk(folder_path):
        for file_name in files:
            file_path = os.path.join(root, file_name)
            content = read_file(file_path)
            print("-----------filepath-------------", file_path)
            contents.append(content)
    return contents    







    
file_path = 'C:\LLAMA2Locally\PDF\PO.pptx'
file_content = read_file(file_path)
print(file_content)


    

